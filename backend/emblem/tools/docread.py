"""Read/extract text from files in any form — so Emblem can read documents.

Supports: txt/md/csv/code, Word (.docx), PDF (.pdf), PowerPoint (.pptx), Excel (.xlsx).
Safe tier (read-only). Files can live under any allowed root (documents/workspaces/exports).
"""

from ..kernel import Tier, tool
from . import fsutil

TEXT_EXT  = {"txt", "md", "csv", "tsv", "json", "py", "js", "ts", "tsx", "jsx", "html",
             "css", "svelte", "sh", "ps1", "sql", "yaml", "yml", "xml", "log"}
IMAGE_EXT = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "tif"}


@tool("docs.read", Tier.SAFE, "Read and extract text from a file (Word/PDF/PowerPoint/Excel/text/code)",
      summarize=lambda a: f"Read {a.get('root', 'documents')}/{a.get('path', '')}")
def read_document(path="", root="documents", max_chars=14000):
    target = fsutil.resolve(root, path)
    if not target.exists():
        return {"ok": False, "error": "File not found."}
    ext = target.suffix.lower().lstrip(".")
    try:
        if ext in TEXT_EXT:
            text = target.read_text(encoding="utf-8", errors="replace")
        elif ext in ("docx", "doc"):
            from docx import Document
            d = Document(str(target))
            text = "\n".join(p.text for p in d.paragraphs if p.text)
        elif ext == "pdf":
            from pypdf import PdfReader
            rd = PdfReader(str(target))
            text = "\n".join((pg.extract_text() or "") for pg in rd.pages)
        elif ext == "pptx":
            from pptx import Presentation
            pr = Presentation(str(target))
            parts = []
            for i, s in enumerate(pr.slides, 1):
                parts.append(f"--- Slide {i} ---")
                for sh in s.shapes:
                    if sh.has_text_frame and sh.text_frame.text:
                        parts.append(sh.text_frame.text)
            text = "\n".join(parts)
        elif ext in ("xlsx", "xlsm"):
            from openpyxl import load_workbook
            wb = load_workbook(str(target), data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join("" if c is None else str(c) for c in row))
            text = "\n".join(parts)
        elif ext in IMAGE_EXT:
            from ..agent.brain import chat_vision
            result = chat_vision(
                "Describe this image in full detail — include any text, charts, objects, people, and context.",
                target,
            )
            text = result.get("reply", "") if result.get("ok") else f"[image analysis unavailable: {result.get('error', 'no vision key')}]"
        else:
            text = target.read_text(encoding="utf-8", errors="replace")
        return {"ok": True, "root": root, "path": path, "chars": len(text), "text": text[:int(max_chars)]}
    except Exception as exc:
        return {"ok": False, "error": str(exc)}
