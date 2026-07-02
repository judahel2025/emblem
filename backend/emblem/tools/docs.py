"""Document generation — turn Markdown content into real, well-formatted files
(.md, .txt, .docx, .pdf, .pptx) saved into the documents folder so they show up in
Files and open in local apps (Word, etc.).

A proper Markdown parser maps headings, paragraphs, bullet/numbered lists, and GitHub
pipe-tables into native Word/PDF/PowerPoint structures. Inline **bold**, *italic*, and
`code` are rendered as real formatting — never left as literal asterisks. The goal: no
scattered tables, no stray markup, clean headings everywhere.

Caution tier: writes a file.
"""

import re
import time

from ..kernel import Tier, tool
from . import fsutil


def _slug(title):
    s = "".join(c.lower() if c.isalnum() else "-" for c in (title or "document"))
    while "--" in s:
        s = s.replace("--", "-")
    return s.strip("-") or "document"


# --- Markdown -> block model -------------------------------------------------

_TABLE_DIVIDER = re.compile(r"^\s*\|?\s*:?-{2,}:?\s*(\|\s*:?-{2,}:?\s*)+\|?\s*$")


def _split_row(line):
    return [c.strip() for c in line.strip().strip("|").split("|")]


def parse_blocks(content):
    """Parse Markdown into a list of typed blocks for the renderers.

    Block shapes:
      ("heading", level, text)
      ("para", text)
      ("bullet", text)         ("number", text)
      ("table", [headers], [[row cells], ...])
      ("code", text)
      ("rule", None)
    """
    lines = (content or "").replace("\r\n", "\n").split("\n")
    blocks = []
    i, n = 0, len(lines)
    para = []

    def flush_para():
        if para:
            blocks.append(("para", " ".join(para).strip()))
            para.clear()

    while i < n:
        line = lines[i]
        s = line.strip()

        # fenced code block
        if s.startswith("```"):
            flush_para()
            i += 1
            buf = []
            while i < n and not lines[i].strip().startswith("```"):
                buf.append(lines[i])
                i += 1
            i += 1  # skip closing fence
            blocks.append(("code", "\n".join(buf)))
            continue

        # table: a header row of pipes followed by a |---|---| divider
        if "|" in s and i + 1 < n and _TABLE_DIVIDER.match(lines[i + 1]):
            flush_para()
            headers = _split_row(s)
            i += 2
            rows = []
            while i < n and "|" in lines[i] and lines[i].strip():
                if not _TABLE_DIVIDER.match(lines[i]):
                    rows.append(_split_row(lines[i]))
                i += 1
            blocks.append(("table", headers, rows))
            continue

        if not s:
            flush_para()
            i += 1
            continue

        if re.match(r"^(-{3,}|\*{3,}|_{3,})$", s):
            flush_para()
            blocks.append(("rule", None))
            i += 1
            continue

        m = re.match(r"^(#{1,6})\s+(.*)$", s)
        if m:
            flush_para()
            blocks.append(("heading", len(m.group(1)), m.group(2).strip()))
            i += 1
            continue

        m = re.match(r"^[-*+]\s+(.*)$", s)
        if m:
            flush_para()
            blocks.append(("bullet", m.group(1).strip()))
            i += 1
            continue

        m = re.match(r"^\d+[.)]\s+(.*)$", s)
        if m:
            flush_para()
            blocks.append(("number", m.group(1).strip()))
            i += 1
            continue

        para.append(s)
        i += 1

    flush_para()
    return blocks


# inline tokens: (kind, text) where kind in {text, bold, italic, code}
_INLINE = re.compile(r"(\*\*[^*]+\*\*|__[^_]+__|\*[^*]+\*|_[^_]+_|`[^`]+`)")


def parse_inline(text):
    out = []
    for part in _INLINE.split(text or ""):
        if not part:
            continue
        if (part.startswith("**") and part.endswith("**")) or (part.startswith("__") and part.endswith("__")):
            out.append(("bold", part[2:-2]))
        elif part.startswith("`") and part.endswith("`"):
            out.append(("code", part[1:-1]))
        elif (part.startswith("*") and part.endswith("*")) or (part.startswith("_") and part.endswith("_")):
            out.append(("italic", part[1:-1]))
        else:
            out.append(("text", part))
    return out


def _plain(text):
    """Strip inline markdown markers to clean text (for pptx / fallbacks)."""
    return "".join(seg for _, seg in parse_inline(text))


# --- DOCX --------------------------------------------------------------------

def _docx_runs(paragraph, text):
    for kind, seg in parse_inline(text):
        run = paragraph.add_run(seg)
        if kind == "bold":
            run.bold = True
        elif kind == "italic":
            run.italic = True
        elif kind == "code":
            run.font.name = "Consolas"


def _make_docx(title, content, target):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    for block in parse_blocks(content):
        kind = block[0]
        if kind == "heading":
            doc.add_heading(_plain(block[2]), level=min(block[1], 4))
        elif kind == "para":
            _docx_runs(doc.add_paragraph(), block[1])
        elif kind == "bullet":
            _docx_runs(doc.add_paragraph(style="List Bullet"), block[1])
        elif kind == "number":
            _docx_runs(doc.add_paragraph(style="List Number"), block[1])
        elif kind == "code":
            p = doc.add_paragraph()
            r = p.add_run(block[1])
            r.font.name = "Consolas"
        elif kind == "rule":
            doc.add_paragraph("_" * 30).alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif kind == "table":
            headers, rows = block[1], block[2]
            cols = max(len(headers), *(len(r) for r in rows)) if rows else len(headers)
            table = doc.add_table(rows=1, cols=cols)
            try:
                table.style = "Light Grid Accent 1"
            except Exception:
                table.style = "Table Grid"
            for c in range(cols):
                cell = table.rows[0].cells[c]
                cell.paragraphs[0].add_run(headers[c] if c < len(headers) else "").bold = True
            for r in rows:
                cells = table.add_row().cells
                for c in range(cols):
                    _docx_runs(cells[c].paragraphs[0], r[c] if c < len(r) else "")
    doc.save(str(target))


# --- PDF ---------------------------------------------------------------------

def _pdf_markup(text):
    """Convert inline markdown to ReportLab's mini-HTML markup."""
    out = []
    for kind, seg in parse_inline(text):
        seg = seg.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if kind == "bold":
            out.append(f"<b>{seg}</b>")
        elif kind == "italic":
            out.append(f"<i>{seg}</i>")
        elif kind == "code":
            out.append(f'<font face="Courier">{seg}</font>')
        else:
            out.append(seg)
    return "".join(out)


def _make_pdf(title, content, target):
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, ListFlowable, ListItem)

    styles = getSampleStyleSheet()
    flow = []
    if title:
        flow.append(Paragraph(title, styles["Title"]))
        flow.append(Spacer(1, 12))

    head = {1: "Heading1", 2: "Heading2", 3: "Heading3", 4: "Heading4", 5: "Heading5", 6: "Heading6"}
    for block in parse_blocks(content):
        kind = block[0]
        if kind == "heading":
            flow.append(Paragraph(_pdf_markup(block[2]), styles[head.get(block[1], "Heading4")]))
        elif kind == "para":
            flow.append(Paragraph(_pdf_markup(block[1]), styles["BodyText"]))
        elif kind in ("bullet", "number"):
            flow.append(ListFlowable(
                [ListItem(Paragraph(_pdf_markup(block[1]), styles["BodyText"]))],
                bulletType="bullet" if kind == "bullet" else "1"))
        elif kind == "code":
            flow.append(Paragraph(block[1].replace(" ", "&nbsp;").replace("\n", "<br/>"), styles["Code"]))
        elif kind == "rule":
            flow.append(Spacer(1, 6))
        elif kind == "table":
            headers, rows = block[1], block[2]
            cols = max(len(headers), *(len(r) for r in rows)) if rows else len(headers)
            data = [[Paragraph(_pdf_markup(headers[c] if c < len(headers) else ""), styles["BodyText"]) for c in range(cols)]]
            for r in rows:
                data.append([Paragraph(_pdf_markup(r[c] if c < len(r) else ""), styles["BodyText"]) for c in range(cols)])
            tbl = Table(data, hAlign="LEFT")
            tbl.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#b8b0a0")),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#efe9dd")),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]))
            flow.append(tbl)
        flow.append(Spacer(1, 6))
    SimpleDocTemplate(str(target), pagesize=LETTER).build(flow)


# --- PPTX --------------------------------------------------------------------

def _make_pptx(title, content, target):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title or "Presentation"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = "Emblem"

    body = None  # text_frame of current content slide

    def new_slide(heading):
        nonlocal body
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = _plain(heading)
        body = slide.placeholders[1].text_frame
        body.clear()
        body._first_used = False
        return slide

    def add_line(text, level=0):
        nonlocal body
        if body is None:
            new_slide("Overview")
        if not getattr(body, "_first_used", False) and body.paragraphs[0].text == "":
            p = body.paragraphs[0]
            body._first_used = True
        else:
            p = body.add_paragraph()
        p.text = _plain(text)
        p.level = level

    for block in parse_blocks(content):
        kind = block[0]
        if kind == "heading":
            new_slide(block[2])
        elif kind == "para":
            add_line(block[1], 0)
        elif kind == "bullet":
            add_line(block[1], 1)
        elif kind == "number":
            add_line(block[1], 1)
        elif kind == "code":
            for ln in block[1].split("\n"):
                add_line(ln, 1)
        elif kind == "table":
            headers, rows = block[1], block[2]
            data = [headers] + rows
            r, c = len(data), max(len(headers), *(len(x) for x in rows)) if rows else len(headers)
            if body is None:
                new_slide("Table")
            slide = prs.slides[-1]
            tbl = slide.shapes.add_table(r, c, Inches(0.6), Inches(2.0), Inches(9), Inches(0.4 * r)).table
            for ci in range(c):
                tbl.cell(0, ci).text = _plain(headers[ci]) if ci < len(headers) else ""
            for ri, row in enumerate(rows, start=1):
                for ci in range(c):
                    tbl.cell(ri, ci).text = _plain(row[ci]) if ci < len(row) else ""
    prs.save(str(target))


@tool("docs.write", Tier.CAUTION, "Generate and save a document (md/txt/docx/pdf/pptx)",
      summarize=lambda a: f"Save {a.get('fmt', 'md')}: {a.get('title', 'document')}",
      target=lambda a: f"documents/{a.get('title', 'document')}.{a.get('fmt', 'md')}")
def write_doc(title="document", content="", fmt="md"):
    fmt = (fmt or "md").lower().lstrip(".")
    name = f"{_slug(title)}-{int(time.time())}.{fmt}"
    target = fsutil.resolve("documents", name)
    target.parent.mkdir(parents=True, exist_ok=True)
    try:
        if fmt == "docx":
            _make_docx(title, content, target)
        elif fmt == "pdf":
            _make_pdf(title, content, target)
        elif fmt == "pptx":
            _make_pptx(title, content, target)
        else:
            if fmt not in ("md", "txt"):
                fmt = "md"; name = f"{_slug(title)}-{int(time.time())}.md"; target = fsutil.resolve("documents", name)
            target.write_text(content or "", encoding="utf-8")
        return {"ok": True, "root": "documents", "path": name, "fmt": fmt}
    except Exception as exc:
        return {"ok": False, "error": str(exc)}
