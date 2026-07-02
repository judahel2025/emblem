"""Analyze tool — AI-powered analysis of images and any document.

analyze.image   — send an image file to Gemini Vision and get a description / opinion.
analyze.ask     — read a doc (any format) then ask the AI a specific question about it.

Both are SAFE tier (read-only, network to Gemini only).
"""

from ..kernel import Tier, tool
from . import fsutil

IMAGE_EXT = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "tif", "svg"}


@tool(
    "analyze.image",
    Tier.SAFE,
    "Analyze an image with AI vision — describe contents, read text, interpret charts",
    network=True,
    summarize=lambda a: f"Analyze image: {a.get('root', 'documents')}/{a.get('path', '')}",
)
def analyze_image(path="", root="documents",
                  prompt="Describe this image in detail. Note any text, charts, people, objects, or relevant context."):
    """Analyse an image file using Gemini Vision.

    Supports: png, jpg, gif, webp, bmp, tiff.
    Falls back to a text description of the filename if no Gemini key is set.
    """
    from ..agent.brain import chat_vision

    target = fsutil.resolve(root, path)
    if not target.exists():
        return {"ok": False, "error": f"File not found: {root}/{path}"}

    ext = target.suffix.lower().lstrip(".")
    if ext not in IMAGE_EXT:
        return {
            "ok": False,
            "error": f"Not an image file (got .{ext}). Use docs.read for documents.",
        }

    result = chat_vision(prompt, target)
    if not result["ok"]:
        return {"ok": False, "error": result.get("reply") or result.get("error")}
    return {
        "ok":          True,
        "root":        root,
        "path":        path,
        "description": result["reply"],
        "model":       result.get("model", "gemini"),
    }


@tool(
    "analyze.ask",
    Tier.SAFE,
    "Read a document (PDF/Word/Excel/PowerPoint/text) and ask the AI a specific question about it",
    network=True,
    summarize=lambda a: f"Analyze {a.get('root', 'documents')}/{a.get('path', '')}: {a.get('question', '')[:50]}",
)
def analyze_ask(path="", root="documents", question="Summarise this document and highlight the key points.",
                max_chars=14000):
    """Extract text from any supported file, then ask the AI a question about it.

    Supported formats: PDF, DOCX, PPTX, XLSX, and all plain-text types.
    Uses the same extraction logic as docs.read, then sends the content to the brain.
    """
    from ..agent.brain import chat

    target = fsutil.resolve(root, path)
    if not target.exists():
        return {"ok": False, "error": f"File not found: {root}/{path}"}

    # --- extract text ---
    ext  = target.suffix.lower().lstrip(".")
    TEXT = {"txt", "md", "csv", "tsv", "json", "py", "js", "ts", "tsx", "jsx", "html",
            "css", "svelte", "sh", "ps1", "sql", "yaml", "yml", "xml", "log"}
    try:
        if ext in TEXT:
            text = target.read_text(encoding="utf-8", errors="replace")
        elif ext in ("docx", "doc"):
            from docx import Document
            text = "\n".join(p.text for p in Document(str(target)).paragraphs if p.text)
        elif ext == "pdf":
            from pypdf import PdfReader
            text = "\n".join((pg.extract_text() or "") for pg in PdfReader(str(target)).pages)
        elif ext == "pptx":
            from pptx import Presentation
            parts = []
            for i, s in enumerate(Presentation(str(target)).slides, 1):
                parts.append(f"--- Slide {i} ---")
                for sh in s.shapes:
                    if sh.has_text_frame and sh.text_frame.text:
                        parts.append(sh.text_frame.text)
            text = "\n".join(parts)
        elif ext in ("xlsx", "xlsm"):
            from openpyxl import load_workbook
            wb    = load_workbook(str(target), data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join("" if c is None else str(c) for c in row))
            text = "\n".join(parts)
        elif ext in IMAGE_EXT:
            # image — route to vision instead
            return analyze_image(path=path, root=root, prompt=question)
        else:
            text = target.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return {"ok": False, "error": str(exc)}

    text = text[:int(max_chars)]
    if not text.strip():
        return {"ok": False, "error": "Could not extract any text from this file."}

    prompt = (
        f"The following is the full content of the file '{target.name}'.\n\n"
        f"---\n{text}\n---\n\n"
        f"Question: {question}"
    )
    result = chat(prompt, max_tokens=2048)
    return {
        "ok":      result["ok"],
        "root":    root,
        "path":    path,
        "answer":  result.get("reply", ""),
        "chars":   len(text),
        "model":   result.get("model", ""),
        "error":   result.get("error") if not result["ok"] else None,
    }
